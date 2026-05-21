import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Network Protocol Visualizer & Error Detection Playground"
subtitle.text = "Interactive OSI Model Simulation and Error Detection Techniques"

# Slide 2: Project Overview
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Project Overview"
body = slide.shapes.placeholders[1].text_frame
body.text = "• React + Vite frontend"
for point in ["• Zustand state management","• OSI layer visualization","• Error detection algorithms (Parity, Checksum, CRC)","• Live deployment on Vercel"]:
    p = body.add_paragraph()
    p.text = point
    p.level = 0

# Slide 3: Features
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Key Features"
body = slide.shapes.placeholders[1].text_frame
features = [
    "OSI Layer Stack with animated packet flow",
    "Control panel for play/pause/step",
    "Info panel showing packet details",
    "Error Detection Playground – toggle parity, checksum, CRC-8",
    "Inject bit errors and see detection results",
    "Responsive design – works on desktop & mobile"
]
for f in features:
    p = body.add_paragraph()
    p.text = f
    p.level = 0

# Slide 4: Architecture Diagram (Real Image)
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
slide.shapes.title.text = "Architecture Overview"
left = Inches(1.5)
top = Inches(1.5)
width = Inches(7)
height = Inches(4.5)
try:
    slide.shapes.add_picture('architecture.png', left, top, width, height)
except Exception as e:
    # Fallback to rectangle if image not found
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = f"[Architecture Diagram]\nImage not found: {e}"

# Slide 5: Error Detection Techniques
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Error Detection Techniques"
body = slide.shapes.placeholders[1].text_frame
techs = [
    "1‑D Parity (even/odd)",
    "2‑D Parity (row & column)",
    "Checksum (8‑bit sum)",
    "CRC‑8 (generator polynomial 0x07)"
]
for t in techs:
    p = body.add_paragraph()
    p.text = t
    p.level = 0

# Slide 6: Live Links
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Live Links"
body = slide.shapes.placeholders[1].text_frame
links = [
    ("GitHub Repository", "https://github.com/chiiinmay/network-protocol-visualizer"),
    ("Vercel Deployment", "https://network-protocol-visualizer.vercel.app")
]
for name, url in links:
    p = body.add_paragraph()
    p.text = f"{name}: {url}"
    p.level = 0

# Save file
prs.save('Network_Protocol_Visualizer.pptx')
print('Presentation generated')
